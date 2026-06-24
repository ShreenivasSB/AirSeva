"""
generate_ppt.py
Generates AirSeva_1M1B_Submission.pptx using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# -- Colour palette -----------------------------------------------------------
BG_COLOR      = RGBColor(232, 244, 253)
NAVY          = RGBColor(26,  58,  92)
MED_BLUE      = RGBColor(44,  82,  130)
ACCENT_BORDER = RGBColor(99,  179, 237)
BOX_FILL      = RGBColor(219, 234, 254)
FOOTER_BG     = RGBColor(26,  58,  92)
FOOTER_FG     = RGBColor(255, 255, 255)
FOOTER_TEXT   = ("1M1B AI for Sustainability · IBM SkillsBuild + AICTE · "
                 "Shreenivas S B · Dayananda Sagar University, Bangalore")

SLIDE_W       = Inches(13.33)
SLIDE_H       = Inches(7.5)
FOOTER_Y      = Inches(7.1)
FOOTER_HEIGHT = Inches(0.4)
FONT          = "Calibri"


# -- Helpers ------------------------------------------------------------------

def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide):
    txb = slide.shapes.add_textbox(Inches(0), FOOTER_Y, SLIDE_W, FOOTER_HEIGHT)
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = FOOTER_TEXT
    run.font.name  = FONT
    run.font.size  = Pt(9)
    run.font.color.rgb = FOOTER_FG
    txb.fill.solid()
    txb.fill.fore_color.rgb = FOOTER_BG
    return txb


def add_slide_title(slide, title_text, y=Inches(0.3), height=Inches(0.7), font_size=28):
    txb = slide.shapes.add_textbox(Inches(0.4), y, Inches(12.53), height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name  = FONT
    run.font.size  = Pt(font_size)
    run.font.bold  = True
    run.font.color.rgb = NAVY
    return txb


def add_rounded_box(slide, left, top, width, height, lines,
                    font_size=12, text_color=NAVY,
                    fill_color=BOX_FILL, border_color=ACCENT_BORDER,
                    bold_first=False, italic=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(5, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name  = FONT
        run.font.size  = Pt(font_size)
        run.font.color.rgb = text_color
        run.font.bold  = (bold_first and i == 0)
        run.font.italic = italic
    return shape


def add_text_box(slide, left, top, width, height,
                 text, font_size=12, color=MED_BLUE,
                 bold=False, italic=False, align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.italic = italic
    return txb


# -- Slide 1: Title -----------------------------------------------------------

def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)

    # Main title
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(1.2))
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\U0001f32b\ufe0f AirSeva"
    run.font.name  = FONT
    run.font.size  = Pt(44)
    run.font.bold  = True
    run.font.color.rgb = NAVY

    # Subtitle
    txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(0.6))
    tf2  = txb2.text_frame
    tf2.word_wrap = True
    p2   = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Agentic Air Quality Health Advisory System for Indian Communities"
    run2.font.name  = FONT
    run2.font.size  = Pt(20)
    run2.font.color.rgb = MED_BLUE

    # Info lines
    info = [
        "Shreenivas S B",
        "Dayananda Sagar University, Bangalore | MCA - Data Science",
        "1M1B AI for Sustainability . IBM SkillsBuild + AICTE",
    ]
    txb3 = slide.shapes.add_textbox(Inches(0.5), Inches(3.25), Inches(12.33), Inches(1.0))
    tf3  = txb3.text_frame
    tf3.word_wrap = True
    for i, line in enumerate(info):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.name  = FONT
        run.font.size  = Pt(14)
        run.font.color.rgb = MED_BLUE

    # Live app box
    shape = slide.shapes.add_shape(5, Inches(2.5), Inches(5.8), Inches(8.33), Inches(0.65))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BOX_FILL
    shape.line.color.rgb = ACCENT_BORDER
    shape.line.width = Pt(1.5)
    tf4 = shape.text_frame
    tf4.word_wrap = False
    p4  = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = "Live App: https://airseva-4uzac5mbekkwmzvdt4rsux.streamlit.app"
    run4.font.name  = FONT
    run4.font.size  = Pt(11)
    run4.font.color.rgb = NAVY

    add_footer(slide)


# -- Slide 2: SDG Alignment ---------------------------------------------------

def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    add_slide_title(slide, "SDG Alignment")

    sdg_boxes = [
        {
            "left": Inches(0.5), "top": Inches(1.3),
            "lines": [
                "SDG 3 - Good Health & Well-Being",
                "",
                "AirSeva delivers real-time, personalised air quality health advisories "
                "to help Indian communities - especially vulnerable groups - protect "
                "their health during high pollution periods.",
            ],
        },
        {
            "left": Inches(6.93), "top": Inches(1.3),
            "lines": [
                "SDG 11 - Sustainable Cities & Communities",
                "",
                "By covering 26 Indian cities including Tier-2 and Tier-3 locations, "
                "AirSeva promotes equitable access to environmental health information, "
                "supporting more sustainable and resilient urban communities.",
            ],
        },
    ]

    for bx in sdg_boxes:
        shape = slide.shapes.add_shape(5, bx["left"], bx["top"], Inches(5.9), Inches(5.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BOX_FILL
        shape.line.color.rgb = ACCENT_BORDER
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        for i, line in enumerate(bx["lines"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line
            run.font.name  = FONT
            run.font.size  = Pt(14) if i == 0 else Pt(13)
            run.font.color.rgb = NAVY
            run.font.bold = (i == 0)

    add_footer(slide)


# -- Slide 3: Problem Statement -----------------------------------------------

def build_slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    add_slide_title(slide, "Problem Statement")

    # Main box
    main_text = (
        "How might we use AI to provide accessible, real-time, and location-specific "
        "air quality health advisories so that communities across India - especially "
        "vulnerable groups - can take timely precautions and become more resilient to air pollution?"
    )
    add_rounded_box(slide, Inches(0.5), Inches(1.25), Inches(12.33), Inches(1.7),
                    [main_text], font_size=16, italic=True, align=PP_ALIGN.CENTER)

    # Three boxes
    boxes = [
        ("The Gap",
         "300M+ Indians exposed to hazardous AQI levels annually. "
         "Existing tools show raw AQI numbers with no health guidance."),
        ("The Vulnerable",
         "Children, elderly, asthma patients, and outdoor workers face "
         "disproportionate health risk from air pollution."),
        ("The Need",
         "A system that translates raw pollutant data into clear, "
         "personalised, actionable health advice - in real time."),
    ]
    box_w   = Inches(3.9)
    box_gap = Inches(0.26)
    box_y   = Inches(3.25)
    box_h   = Inches(3.45)

    for i, (header, body) in enumerate(boxes):
        left_x = Inches(0.5) + i * (box_w + box_gap)
        shape  = slide.shapes.add_shape(5, left_x, box_y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = BOX_FILL
        shape.line.color.rgb = ACCENT_BORDER
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        for j, line in enumerate(["" if k == 1 else (header if k == 0 else body)
                                   for k in range(3)]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line
            run.font.name  = FONT
            run.font.size  = Pt(14) if j == 0 else Pt(12)
            run.font.color.rgb = NAVY
            run.font.bold = (j == 0)

    add_footer(slide)


# -- Slide 4: AI Solution -----------------------------------------------------

def build_slide4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    add_slide_title(slide, "AI Solution - 4-Agent Pipeline")

    pipeline = [
        ["Agent 1", "Data Fetcher", "WAQI API + Open-Meteo", "Live AQI + 7-day PM2.5 history"],
        ["Agent 2", "ML Predictor", "Random Forest", "Risk: Low / Moderate / High"],
        ["Agent 3", "Advisor", "IBM Granite 4", "WatsonX Frankfurt"],
        ["Agent 4", "Reporter", "WHO 2021 checks", "PDF report"],
        ["Output", "Streamlit Dashboard", "AQI gauge . Charts", "Downloadable PDF"],
    ]
    emojis = ["\U0001f4e1", "\U0001f9e0", "\U0001f535", "\U0001f4cb", "\U0001f4ca"]

    box_w   = Inches(2.2)
    box_h   = Inches(3.8)
    arrow_w = Inches(0.35)
    total_w = 5 * box_w + 4 * arrow_w
    start_x = (SLIDE_W - total_w) / 2
    box_y   = Inches(1.4)

    for i, (lines, emoji) in enumerate(zip(pipeline, emojis)):
        left_x = start_x + i * (box_w + arrow_w)
        shape  = slide.shapes.add_shape(5, left_x, box_y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = BOX_FILL
        shape.line.color.rgb = ACCENT_BORDER
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        all_lines = [emoji + " " + lines[0]] + lines[1:]
        for j, line in enumerate(all_lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.name  = FONT
            run.font.size  = Pt(11)
            run.font.color.rgb = NAVY
            run.font.bold = (j == 0)

        if i < 4:
            arrow_x = left_x + box_w
            arrow_y = box_y + box_h / 2 - Inches(0.15)
            arrow   = slide.shapes.add_shape(13, arrow_x, arrow_y, arrow_w, Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_BORDER
            arrow.line.color.rgb = ACCENT_BORDER

    stack_text = ("Stack: Python . Streamlit . scikit-learn . "
                  "IBM Granite 4 (WatsonX) . WAQI API . WHO 2021 Guidelines")
    add_text_box(slide, Inches(0.5), Inches(5.45), Inches(12.33), Inches(0.4),
                 stack_text, font_size=11, color=MED_BLUE)
    add_footer(slide)


# -- Slide 5: Target Users ----------------------------------------------------

def build_slide5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    add_slide_title(slide, "Target Users")

    users = [
        ("\U0001f465 General Public",
         "Citizens seeking real-time AQI information and daily outdoor activity guidance."),
        ("\U0001f3e5 Vulnerable Groups",
         "Children, elderly, pregnant women, asthma/COPD patients needing cautious, tailored advisories."),
        ("\U0001f3d8\ufe0f ASHA Workers & Health Volunteers",
         "Frontline workers using AirSeva to counsel vulnerable households and run awareness campaigns."),
        ("\U0001f3db\ufe0f Urban Local Bodies & Policy Makers",
         "Municipal teams referencing city-level AQI trends to inform public health advisories and policy."),
    ]
    positions = [
        (Inches(0.5), Inches(1.3)),
        (Inches(6.93), Inches(1.3)),
        (Inches(0.5), Inches(4.15)),
        (Inches(6.93), Inches(4.15)),
    ]
    for (header, body), (left_x, top_y) in zip(users, positions):
        shape = slide.shapes.add_shape(5, left_x, top_y, Inches(5.9), Inches(2.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BOX_FILL
        shape.line.color.rgb = ACCENT_BORDER
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        for i, line in enumerate([header, "", body]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line
            run.font.name  = FONT
            run.font.size  = Pt(12)
            run.font.color.rgb = NAVY
            run.font.bold = (i == 0)

    add_footer(slide)


# -- Slide 6: Responsible AI --------------------------------------------------

def build_slide6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    add_slide_title(slide, "Responsible AI Considerations")

    rai = [
        ("\u2696\ufe0f Fairness",
         "WHO 2021 standards applied equally across all 26 cities - "
         "metro, Tier-2, and Tier-3. No metro-centric bias."),
        ("\U0001f50d Transparency",
         "All agent outputs (pollutant data, ML prediction, advisory) "
         "are shown explicitly. No black-box responses."),
        ("\U0001f6e1\ufe0f Ethics",
         "Medical disclaimer displayed prominently. Advice framed as "
         "precautionary guidance, not medical prescription."),
        ("\U0001f512 Privacy",
         "No personal data collected or stored. All data is public "
         "city-level AQI from open APIs. No user tracking."),
    ]
    positions = [
        (Inches(0.5), Inches(1.3)),
        (Inches(6.93), Inches(1.3)),
        (Inches(0.5), Inches(4.15)),
        (Inches(6.93), Inches(4.15)),
    ]
    for (header, body), (left_x, top_y) in zip(rai, positions):
        shape = slide.shapes.add_shape(5, left_x, top_y, Inches(5.9), Inches(2.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BOX_FILL
        shape.line.color.rgb = ACCENT_BORDER
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        for i, line in enumerate([header, "", body]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line
            run.font.name  = FONT
            run.font.size  = Pt(12)
            run.font.color.rgb = NAVY
            run.font.bold = (i == 0)

    add_footer(slide)


# -- Slide 7: Expected Impact -------------------------------------------------

def build_slide7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    add_slide_title(slide, "Expected Impact")

    impact = [
        ("\U0001f30d Scale",
         ["26 Indian cities covered", "Metro + Tier-2 + Tier-3", "Equitable health access"]),
        ("\U0001f464 Individual",
         ["Personalised vulnerability scoring", "Actionable daily guidance", "Empowers health decisions"]),
        ("\U0001f3d9\ufe0f Community",
         ["Supports ASHA worker outreach", "Informs municipal policy", "Builds pollution resilience"]),
    ]
    box_w   = Inches(3.9)
    box_h   = Inches(4.5)
    box_gap = Inches(0.26)
    box_y   = Inches(1.3)

    for i, (header, items) in enumerate(impact):
        left_x = Inches(0.5) + i * (box_w + box_gap)
        shape  = slide.shapes.add_shape(5, left_x, box_y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = BOX_FILL
        shape.line.color.rgb = ACCENT_BORDER
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        all_lines = [header, ""] + items
        for j, line in enumerate(all_lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.name  = FONT
            run.font.size  = Pt(14) if j == 0 else Pt(12)
            run.font.color.rgb = NAVY
            run.font.bold = (j == 0)

    add_text_box(slide, Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.45),
                 "Aligned with India's national health and sustainability goals under SDG 3 and SDG 11",
                 font_size=12, color=NAVY, italic=True)
    add_footer(slide)


# -- Slide 8: Prototype & Demo ------------------------------------------------

def build_slide8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    add_slide_title(slide, "Prototype & Demo")

    # Left - Agent Flow
    flow = [
        "Agent Flow", "",
        "User selects city + health profile", "\u2193",
        "Agent 1: Fetches live AQI + pollutant data", "\u2193",
        "Agent 2: Random Forest predicts risk level", "\u2193",
        "Agent 3: IBM Granite 4 generates advisory", "\u2193",
        "Agent 4: Compiles PDF report", "\u2193",
        "Dashboard displays results + download",
    ]
    shape_l = slide.shapes.add_shape(5, Inches(0.5), Inches(1.3), Inches(5.9), Inches(5.55))
    shape_l.fill.solid()
    shape_l.fill.fore_color.rgb = BOX_FILL
    shape_l.line.color.rgb = ACCENT_BORDER
    shape_l.line.width = Pt(1.5)
    tf_l = shape_l.text_frame
    tf_l.word_wrap = True
    for i, line in enumerate(flow):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.name  = FONT
        run.font.size  = Pt(12)
        run.font.color.rgb = NAVY
        run.font.bold = (i == 0)

    # Right - Key Features
    feats = [
        "Key Features", "",
        "\u2705 GPS location detection",
        "\u2705 26 Indian cities supported",
        "\u2705 WHO 2021 guideline comparison",
        "\u2705 Personalised vulnerability score (0-8)",
        "\u2705 AQI gauge + pollutant bar chart",
        "\u2705 IBM Granite 4 health advisory",
        "\u2705 Downloadable PDF report",
        "\u2705 Medical disclaimer included",
    ]
    shape_r = slide.shapes.add_shape(5, Inches(6.93), Inches(1.3), Inches(5.9), Inches(5.55))
    shape_r.fill.solid()
    shape_r.fill.fore_color.rgb = BOX_FILL
    shape_r.line.color.rgb = ACCENT_BORDER
    shape_r.line.width = Pt(1.5)
    tf_r = shape_r.text_frame
    tf_r.word_wrap = True
    for i, line in enumerate(feats):
        p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name  = FONT
        run.font.size  = Pt(12)
        run.font.color.rgb = NAVY
        run.font.bold = (i == 0)

    # Links
    links_txb = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.5))
    tf_links  = links_txb.text_frame
    tf_links.word_wrap = True
    for i, line in enumerate([
        "\U0001f517 Live App: https://airseva-4uzac5mbekkwmzvdt4rsux.streamlit.app",
        "\U0001f4c1 GitHub: https://github.com/ShreenivasSB/AirSeva",
    ]):
        p = tf_links.paragraphs[0] if i == 0 else tf_links.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.name  = FONT
        run.font.size  = Pt(12)
        run.font.color.rgb = MED_BLUE

    add_footer(slide)


# -- Slide 9: Impact Statement ------------------------------------------------

def build_slide9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    add_slide_title(slide, "Impact Statement")

    # Large quote box
    quote = (
        "AirSeva empowers individuals and communities to make informed, health-conscious "
        "decisions based on real-time air quality data across 26 Indian cities. By translating "
        "raw pollutant data into clear, actionable advisories, it bridges the gap between "
        "complex environmental monitoring and everyday public health behaviour."
    )
    shape_q = slide.shapes.add_shape(5, Inches(0.5), Inches(1.25), Inches(12.33), Inches(2.4))
    shape_q.fill.solid()
    shape_q.fill.fore_color.rgb = BOX_FILL
    shape_q.line.color.rgb = ACCENT_BORDER
    shape_q.line.width = Pt(1.5)
    tf_q = shape_q.text_frame
    tf_q.word_wrap = True
    p_q  = tf_q.paragraphs[0]
    p_q.alignment = PP_ALIGN.CENTER
    run_q = p_q.add_run()
    run_q.text = quote
    run_q.font.name  = FONT
    run_q.font.size  = Pt(15)
    run_q.font.italic = True
    run_q.font.color.rgb = NAVY

    # Two side boxes
    side = [
        ("Who Benefits",
         "Citizens . Vulnerable groups . ASHA workers . Urban local bodies"),
        ("What Changes",
         "Reduced gap between AQI data and public action . "
         "Proactive health behaviour . "
         "Supports India's SDG 3 + SDG 11 goals"),
    ]
    for i, (header, body) in enumerate(side):
        left_x = Inches(0.5) if i == 0 else Inches(6.93)
        shape  = slide.shapes.add_shape(5, left_x, Inches(3.9), Inches(5.9), Inches(2.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BOX_FILL
        shape.line.color.rgb = ACCENT_BORDER
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        for j, line in enumerate([header, "", body]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line
            run.font.name  = FONT
            run.font.size  = Pt(14) if j == 0 else Pt(12)
            run.font.color.rgb = NAVY
            run.font.bold = (j == 0)

    add_text_box(slide, Inches(0.5), Inches(6.55), Inches(12.33), Inches(0.4),
                 "Built with IBM Granite 4 . WatsonX AI . Random Forest ML . WAQI API . WHO 2021 Standards",
                 font_size=11, color=MED_BLUE)
    add_footer(slide)


# -- Main ---------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)
    build_slide4(prs)
    build_slide5(prs)
    build_slide6(prs)
    build_slide7(prs)
    build_slide8(prs)
    build_slide9(prs)

    output_path = "AirSeva_1M1B_Submission.pptx"
    prs.save(output_path)
    print(f"PPT generated successfully: {output_path}")


if __name__ == "__main__":
    main()
