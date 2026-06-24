"""
redesign_ppt.py
Rebuilds AirSeva_1M1B_Submission_Redesigned.pptx with a clean,
professional layout. Text content unchanged; only visual design changes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# Screenshot of the live AirSeva Streamlit app (captured via browser)
APP_SCREENSHOT = (
    r"C:\Users\shree\.gemini\antigravity-ide\brain"
    r"\3782a8ba-ef8c-44fb-8f0a-50f73dec9bb5"
    r"\airseva_dashboard_1782325407729.png"
)

# ── Palette ─────────────────────────────────────────────────────────────────
WHITE       = RGBColor(255, 255, 255)
NAVY        = RGBColor(15,  40,  80)
DARK_GREY   = RGBColor(50,  50,  50)
MID_GREY    = RGBColor(100, 100, 100)
LIGHT_GREY  = RGBColor(80,  80,  80)
BOX_FILL    = RGBColor(240, 245, 252)

SLIDE_W     = Inches(13.33)
SLIDE_H     = Inches(7.5)
FONT        = "Calibri"

FOOTER_TEXT = ("1M1B AI for Sustainability  \u00b7  IBM SkillsBuild + AICTE  "
               "\u00b7  Shreenivas S B  \u00b7  Dayananda Sagar University, Bangalore")

# ── Low-level helpers ────────────────────────────────────────────────────────

def white_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def remove_shadow_xml(shape):
    """Strip effectLst from spPr so no shadow is rendered."""
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    for tag in [qn("a:effectLst"), qn("a:effectDag")]:
        el = spPr.find(tag)
        if el is not None:
            spPr.remove(el)


def no_line(shape):
    """Remove visible border from shape."""
    shape.line.fill.background()


def flat_box(slide, left, top, width, height):
    """Add a flat rectangle: BOX_FILL, no border, no shadow."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BOX_FILL
    no_line(shape)
    remove_shadow_xml(shape)
    return shape


def divider(slide):
    """Thin navy horizontal rule under the slide title."""
    rule = slide.shapes.add_shape(1,
        Inches(0.4), Inches(0.88), Inches(12.5), Inches(0.025))
    rule.fill.solid()
    rule.fill.fore_color.rgb = NAVY
    no_line(rule)
    remove_shadow_xml(rule)
    return rule


def footer(slide):
    """Full-width navy footer bar."""
    bar = slide.shapes.add_shape(1,
        Inches(0), Inches(7.05), SLIDE_W, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    no_line(bar)
    remove_shadow_xml(bar)
    tf = bar.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = FOOTER_TEXT
    run.font.name  = FONT
    run.font.size  = Pt(9)
    run.font.color.rgb = WHITE
    run.font.bold  = False
    return bar


def slide_title(slide, text):
    """Bold navy slide title at top-left."""
    txb = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.6))
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = Pt(36)
    run.font.bold  = True
    run.font.color.rgb = NAVY
    return txb


def add_tb(slide, left, top, width, height):
    """Plain textbox helper."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.text_frame.word_wrap = True
    return txb


def set_para(para, text, size, color, bold=False, italic=False,
             align=PP_ALIGN.LEFT, space_after=Pt(4)):
    para.alignment = align
    para.space_after = space_after
    run = para.add_run()
    run.text = text
    run.font.name   = FONT
    run.font.size   = size
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run


def box_content(shape, heading, body_lines,
                h_size=Pt(16), b_size=Pt(15),
                pad_l=Inches(0.25), pad_t=Inches(0.25)):
    """Fill a shape's text frame with a bold heading + body lines."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left  = pad_l
    tf.margin_right = pad_l
    tf.margin_top   = pad_t
    tf.margin_bottom = pad_t

    # heading
    p0 = tf.paragraphs[0]
    set_para(p0, heading, h_size, NAVY, bold=True)

    # blank spacer
    pb = tf.add_paragraph()
    set_para(pb, "", Pt(6), DARK_GREY, space_after=Pt(2))

    # body
    for i, line in enumerate(body_lines):
        p = tf.add_paragraph()
        set_para(p, line, b_size, DARK_GREY,
                 space_after=Pt(5 if i < len(body_lines)-1 else 0))


def arrow_shape(slide, left, top):
    """Small navy filled rectangle acting as a right-pointing arrow."""
    arr = slide.shapes.add_shape(1, left, top, Inches(0.3), Inches(0.12))
    arr.fill.solid()
    arr.fill.fore_color.rgb = NAVY
    no_line(arr)
    remove_shadow_xml(arr)
    return arr


# ── Slide builders ───────────────────────────────────────────────────────────

def slide1(prs):
    """Title slide."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(sl)

    # Main title "AirSeva"
    t = add_tb(sl, Inches(0.4), Inches(1.8), Inches(12.5), Inches(0.9))
    p = t.text_frame.paragraphs[0]
    set_para(p, "AirSeva", Pt(54), NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Subtitle
    t2 = add_tb(sl, Inches(0.4), Inches(2.7), Inches(12.5), Inches(0.6))
    p2 = t2.text_frame.paragraphs[0]
    set_para(p2, "Agentic Air Quality Health Advisory System for Indian Communities",
             Pt(22), LIGHT_GREY, align=PP_ALIGN.CENTER)

    # Info block
    info_lines = [
        "Shreenivas S B",
        "Dayananda Sagar University, Bangalore  |  MCA \u2013 Data Science",
        "1M1B AI for Sustainability  \u00b7  IBM SkillsBuild + AICTE",
    ]
    t3 = add_tb(sl, Inches(0.4), Inches(3.5), Inches(12.5), Inches(0.9))
    tf3 = t3.text_frame
    tf3.word_wrap = True
    for i, line in enumerate(info_lines):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        set_para(p, line, Pt(15), MID_GREY, align=PP_ALIGN.CENTER, space_after=Pt(3))

    # Live URL box
    url_box = flat_box(sl, Inches(2.5), Inches(5.0), Inches(8.3), Inches(0.55))
    tf4 = url_box.text_frame
    tf4.word_wrap = False
    tf4.margin_left  = Inches(0.2)
    tf4.margin_right = Inches(0.2)
    tf4.margin_top   = Inches(0.1)
    p4 = tf4.paragraphs[0]
    set_para(p4, "Live App: https://airseva-4uzac5mbekkwmzvdt4rsux.streamlit.app",
             Pt(13), NAVY, align=PP_ALIGN.CENTER)

    footer(sl)


def slide2(prs):
    """SDG Alignment."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(sl)
    slide_title(sl, "SDG Alignment")
    divider(sl)

    sdg = [
        ("SDG 3 \u2014 Good Health & Well-Being",
         ["AirSeva delivers real-time, personalised air quality health advisories "
          "to help Indian communities \u2014 especially vulnerable groups \u2014 protect "
          "their health during high pollution periods."]),
        ("SDG 11 \u2014 Sustainable Cities & Communities",
         ["By covering 26 Indian cities including Tier-2 and Tier-3 locations, "
          "AirSeva promotes equitable access to environmental health information, "
          "supporting more sustainable and resilient urban communities."]),
    ]
    lefts = [Inches(0.4), Inches(6.9)]
    for (heading, body), lft in zip(sdg, lefts):
        box = flat_box(sl, lft, Inches(1.1), Inches(5.9), Inches(5.6))
        box_content(box, heading, body)

    footer(sl)


def slide3(prs):
    """Problem Statement."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(sl)
    slide_title(sl, "Problem Statement")
    divider(sl)

    # Wide HMW box
    hmw_box = flat_box(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(1.5))
    tf = hmw_box.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top   = Inches(0.2)
    p = tf.paragraphs[0]
    set_para(p,
        "How might we use AI to provide accessible, real-time, and location-specific "
        "air quality health advisories so that communities across India \u2014 especially "
        "vulnerable groups \u2014 can take timely precautions and become more resilient to air pollution?",
        Pt(16), DARK_GREY, italic=True, align=PP_ALIGN.CENTER)

    # Three boxes
    three = [
        ("The Gap",
         ["300M+ Indians exposed to hazardous AQI levels annually.",
          "Existing tools show raw AQI numbers with no health guidance."]),
        ("The Vulnerable",
         ["Children, elderly, asthma patients, and outdoor workers face "
          "disproportionate health risk from air pollution."]),
        ("The Need",
         ["A system that translates raw pollutant data into clear, "
          "personalised, actionable health advice \u2014 in real time."]),
    ]
    lefts = [Inches(0.4), Inches(4.65), Inches(8.9)]
    for (heading, body), lft in zip(three, lefts):
        box = flat_box(sl, lft, Inches(2.85), Inches(3.9), Inches(3.85))
        box_content(box, heading, body)

    footer(sl)


def slide4(prs):
    """AI Solution \u2013 4-Agent Pipeline."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(sl)
    slide_title(sl, "AI Solution \u2013 4-Agent Pipeline")
    divider(sl)

    agents = [
        ("Agent 1", ["Data Fetcher", "WAQI API + Open-Meteo", "Live AQI + 7-day PM2.5 history"]),
        ("Agent 2", ["ML Predictor", "Random Forest", "Risk: Low / Moderate / High"]),
        ("Agent 3", ["Advisor", "IBM Granite 4", "WatsonX Frankfurt"]),
        ("Agent 4", ["Reporter", "WHO 2021 checks", "PDF report"]),
        ("Output",  ["Streamlit Dashboard", "AQI gauge + Charts", "Downloadable PDF"]),
    ]
    box_lefts  = [Inches(0.4), Inches(2.75), Inches(5.1), Inches(7.45), Inches(9.8)]
    arr_lefts  = [Inches(2.52), Inches(4.87), Inches(7.22), Inches(9.57)]
    arr_center = Inches(3.25) - Inches(0.06)   # vertical center of arrow

    for (heading, body), lft in zip(agents, box_lefts):
        box = flat_box(sl, lft, Inches(1.1), Inches(2.1), Inches(4.3))
        box_content(box, heading, body,
                    h_size=Pt(14), b_size=Pt(13))

    for al in arr_lefts:
        arrow_shape(sl, al, arr_center)

    # Stack line
    t = add_tb(sl, Inches(0.4), Inches(5.6), Inches(12.5), Inches(0.35))
    p = t.text_frame.paragraphs[0]
    set_para(p,
        "Stack: Python  \u00b7  Streamlit  \u00b7  scikit-learn  \u00b7  "
        "IBM Granite 4 (WatsonX)  \u00b7  WAQI API  \u00b7  WHO 2021 Guidelines",
        Pt(12), MID_GREY, align=PP_ALIGN.CENTER)

    footer(sl)


def slide5(prs):
    """Target Users."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(sl)
    slide_title(sl, "Target Users")
    divider(sl)

    users = [
        (Inches(0.4), Inches(1.1),
         "General Public",
         ["Citizens seeking real-time AQI information and daily outdoor activity guidance."]),
        (Inches(6.9), Inches(1.1),
         "Vulnerable Groups",
         ["Children, elderly, pregnant women, asthma/COPD patients needing cautious, tailored advisories."]),
        (Inches(0.4), Inches(4.0),
         "ASHA Workers & Health Volunteers",
         ["Frontline workers using AirSeva to counsel vulnerable households and run awareness campaigns."]),
        (Inches(6.9), Inches(4.0),
         "Urban Local Bodies & Policy Makers",
         ["Municipal teams referencing city-level AQI trends to inform public health advisories and policy."]),
    ]
    for lft, tp, heading, body in users:
        box = flat_box(sl, lft, tp, Inches(5.9), Inches(2.65))
        box_content(box, heading, body)

    footer(sl)


def slide6(prs):
    """Responsible AI Considerations."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(sl)
    slide_title(sl, "Responsible AI Considerations")
    divider(sl)

    rai = [
        (Inches(0.4), Inches(1.1),
         "Fairness",
         ["WHO 2021 standards applied equally across all 26 cities \u2014 "
          "metro, Tier-2, and Tier-3. No metro-centric bias."]),
        (Inches(6.9), Inches(1.1),
         "Transparency",
         ["All agent outputs (pollutant data, ML prediction, advisory) "
          "are shown explicitly. No black-box responses."]),
        (Inches(0.4), Inches(4.0),
         "Ethics",
         ["Medical disclaimer displayed prominently. Advice framed as "
          "precautionary guidance, not medical prescription."]),
        (Inches(6.9), Inches(4.0),
         "Privacy",
         ["No personal data collected or stored. All data is public "
          "city-level AQI from open APIs. No user tracking."]),
    ]
    for lft, tp, heading, body in rai:
        box = flat_box(sl, lft, tp, Inches(5.9), Inches(2.65))
        box_content(box, heading, body)

    footer(sl)


def slide7(prs):
    """Expected Impact."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(sl)
    slide_title(sl, "Expected Impact")
    divider(sl)

    impact = [
        (Inches(0.4),
         "Scale",
         ["26 Indian cities covered", "Metro + Tier-2 + Tier-3", "Equitable health access"]),
        (Inches(4.65),
         "Individual",
         ["Personalised vulnerability scoring", "Actionable daily guidance", "Empowers health decisions"]),
        (Inches(8.9),
         "Community",
         ["Supports ASHA worker outreach", "Informs municipal policy", "Builds pollution resilience"]),
    ]
    for lft, heading, body in impact:
        box = flat_box(sl, lft, Inches(1.1), Inches(3.9), Inches(4.5))
        box_content(box, heading, body)

    # Italic caption
    t = add_tb(sl, Inches(0.4), Inches(5.8), Inches(12.5), Inches(0.4))
    p = t.text_frame.paragraphs[0]
    set_para(p,
        "Aligned with India\u2019s national health and sustainability goals under SDG 3 and SDG 11",
        Pt(13), MID_GREY, italic=True, align=PP_ALIGN.CENTER)

    footer(sl)


def slide8(prs):
    """Prototype & Demo — architecture diagram left, app screenshot right."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(sl)
    slide_title(sl, "Prototype & Demo")
    divider(sl)

    ARCH_PNG = r"C:\SRP_PROJECT\assets\architecture_diagram.png"

    # ── Left panel: Architecture diagram image ──────────────────────────────
    lbl_l = add_tb(sl, Inches(0.4), Inches(1.1), Inches(6.0), Inches(0.32))
    set_para(lbl_l.text_frame.paragraphs[0],
             "System Architecture \u2014 4-Agent Pipeline",
             Pt(13), NAVY, bold=True, align=PP_ALIGN.LEFT)

    arch_left   = Inches(0.4)
    arch_top    = Inches(1.48)
    arch_width  = Inches(6.0)
    arch_height = Inches(5.25)

    frame_l = flat_box(sl, arch_left, arch_top, arch_width, arch_height)
    if os.path.exists(ARCH_PNG):
        sl.shapes.add_picture(
            ARCH_PNG,
            arch_left  + Inches(0.05),
            arch_top   + Inches(0.05),
            arch_width  - Inches(0.1),
            arch_height - Inches(0.1),
        )

    # ── Right panel: Live app screenshot ────────────────────────────────────
    lbl_r = add_tb(sl, Inches(6.7), Inches(1.1), Inches(6.23), Inches(0.32))
    set_para(lbl_r.text_frame.paragraphs[0],
             "Live App \u2014 AirSeva Dashboard",
             Pt(13), NAVY, bold=True, align=PP_ALIGN.LEFT)

    img_left   = Inches(6.7)
    img_top    = Inches(1.48)
    img_width  = Inches(6.23)
    img_height = Inches(5.25)

    frame_r = flat_box(sl, img_left, img_top, img_width, img_height)
    if os.path.exists(APP_SCREENSHOT):
        sl.shapes.add_picture(
            APP_SCREENSHOT,
            img_left  + Inches(0.05),
            img_top   + Inches(0.05),
            img_width  - Inches(0.1),
            img_height - Inches(0.1),
        )

    # ── Link line ───────────────────────────────────────────────────────────
    t = add_tb(sl, Inches(0.4), Inches(6.88), Inches(12.53), Inches(0.28))
    p = t.text_frame.paragraphs[0]
    set_para(p,
        "Live App: https://airseva-4uzac5mbekkwmzvdt4rsux.streamlit.app"
        "   |   GitHub: https://github.com/ShreenivasSB/AirSeva",
        Pt(12), NAVY, align=PP_ALIGN.CENTER)

    footer(sl)


def slide9(prs):
    """Impact Statement."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(sl)
    slide_title(sl, "Impact Statement")
    divider(sl)

    # Wide quote box
    q_box = flat_box(sl, Inches(0.4), Inches(1.1), Inches(12.5), Inches(2.2))
    tf = q_box.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top   = Inches(0.2)
    p = tf.paragraphs[0]
    set_para(p,
        "AirSeva empowers individuals and communities to make informed, health-conscious "
        "decisions based on real-time air quality data across 26 Indian cities. By translating "
        "raw pollutant data into clear, actionable advisories, it bridges the gap between "
        "complex environmental monitoring and everyday public health behaviour.",
        Pt(17), DARK_GREY, italic=True, align=PP_ALIGN.CENTER)

    # Two side boxes
    sides = [
        (Inches(0.4),
         "Who Benefits",
         ["Citizens  \u00b7  Vulnerable groups  \u00b7  ASHA workers  \u00b7  Urban local bodies"]),
        (Inches(6.9),
         "What Changes",
         ["Reduced gap between AQI data and public action  \u00b7  Proactive health behaviour  "
          "\u00b7  Supports India\u2019s SDG 3 + SDG 11 goals"]),
    ]
    for lft, heading, body in sides:
        box = flat_box(sl, lft, Inches(3.55), Inches(5.9), Inches(2.2))
        box_content(box, heading, body)

    # Credits line
    t = add_tb(sl, Inches(0.4), Inches(5.95), Inches(12.5), Inches(0.35))
    p = t.text_frame.paragraphs[0]
    set_para(p,
        "Built with IBM Granite 4  \u00b7  WatsonX AI  \u00b7  Random Forest ML  "
        "\u00b7  WAQI API  \u00b7  WHO 2021 Standards",
        Pt(12), MID_GREY, align=PP_ALIGN.CENTER)

    footer(sl)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    slide7(prs)
    slide8(prs)
    slide9(prs)

    out = "AirSeva_1M1B_Submission_Redesigned.pptx"
    prs.save(out)
    print(f"Redesigned PPT saved as {out}")


if __name__ == "__main__":
    main()
